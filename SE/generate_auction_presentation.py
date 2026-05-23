from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Brand/theme colors
PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)      # Dark blue
SECONDARY = RGBColor(0x4D, 0xA3, 0xFF)    # Light blue
ACCENT = RGBColor(0xF5, 0xA6, 0x23)       # Gold/Orange
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)     # Light grey
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x2A, 0x36)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_standard_header(slide, title, slide_no):
    """Add a consistent professional title/header style for content slides."""
    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.62)
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = BG_LIGHT
    top_band.line.fill.background()

    accent_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.58), SLIDE_WIDTH, Inches(0.04)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = SECONDARY
    accent_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.11), Inches(10.6), Inches(0.38))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Poppins"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    no_box = slide.shapes.add_textbox(Inches(12.2), Inches(0.16), Inches(0.9), Inches(0.3))
    no_tf = no_box.text_frame
    no_tf.clear()
    p2 = no_tf.paragraphs[0]
    p2.text = f"{slide_no}/10"
    p2.font.name = "Calibri"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.RIGHT


def set_box_text(shape, text, font_name="Calibri", size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def add_bullets(slide, x, y, w, h, bullets, font_size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()

    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)


def add_icon_card(slide, x, y, w, h, title, subtitle, badge_text, badge_color):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD9, 0xE2, 0xEC)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.62), Inches(0.62)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    set_box_text(badge, badge_text, font_name="Montserrat", size=12, bold=True, color=WHITE)

    title_box = slide.shapes.add_textbox(Inches(x + 0.95), Inches(y + 0.16), Inches(w - 1.0), Inches(0.36))
    tf = title_box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Poppins"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY

    desc_box = slide.shapes.add_textbox(Inches(x + 0.95), Inches(y + 0.52), Inches(w - 1.05), Inches(h - 0.62))
    desc_tf = desc_box.text_frame
    desc_tf.clear()
    p2 = desc_tf.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Calibri"
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_DARK


def add_arrow(slide, x, y, w=0.45, h=0.22):
    arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = SECONDARY
    arr.line.fill.background()


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(4.2), SLIDE_HEIGHT
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = PRIMARY
    left_panel.line.fill.background()

    accent_strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.85), Inches(4.2), Inches(0.65)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = ACCENT
    accent_strip.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.52), Inches(1.35), Inches(3.2), Inches(2.2))
    tf = title_box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Auction"
    p1.font.name = "Poppins"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "Management"
    p2.font.name = "Poppins"
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    p3 = tf.add_paragraph()
    p3.text = "System"
    p3.font.name = "Poppins"
    p3.font.size = Pt(48)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.1), Inches(7.8), Inches(0.8))
    set_box_text(
        subtitle_box,
        "Software Engineering Project Presentation",
        font_name="Montserrat",
        size=24,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.LEFT,
    )

    details = [
        "Presenter: [Your Name]",
        "Course: [Course Name / Code]",
        "Institution: [Institution Name]",
        f"Date: {date.today().strftime('%B %d, %Y')}",
    ]
    add_bullets(slide, 4.95, 2.2, 7.3, 2.6, details, font_size=20)

    # Simple gavel icon (composed with shapes)
    handle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.15), Inches(4.22), Inches(2.08), Inches(0.22)
    )
    handle.fill.solid()
    handle.fill.fore_color.rgb = PRIMARY
    handle.line.fill.background()
    handle.rotation = -26

    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.35), Inches(3.55), Inches(1.18), Inches(0.42)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = ACCENT
    head.line.fill.background()
    head.rotation = -26

    head2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.0), Inches(3.22), Inches(1.18), Inches(0.42)
    )
    head2.fill.solid()
    head2.fill.fore_color.rgb = ACCENT
    head2.line.fill.background()
    head2.rotation = -26

    base = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.42), Inches(5.15), Inches(2.82), Inches(0.28)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = SECONDARY
    base.line.fill.background()


def add_slide_2_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "2. Introduction", 2)

    bullets = [
        "An Auction Management System is a web-based platform for conducting online auctions.",
        "Sellers can publish auction items with descriptions, pricing, and timelines.",
        "Buyers can place competitive bids in real time from any location.",
        "The system improves transparency, automation, and operational efficiency.",
    ]
    add_bullets(slide, 0.8, 1.2, 7.0, 4.7, bullets, font_size=18)

    add_icon_card(
        slide,
        8.0,
        1.4,
        4.8,
        1.45,
        "Transparent Process",
        "Bids are visible and traceable throughout the auction lifecycle.",
        "TR",
        PRIMARY,
    )
    add_icon_card(
        slide,
        8.0,
        3.0,
        4.8,
        1.45,
        "Automated Operations",
        "Timers, bid updates, and winner selection happen automatically.",
        "AU",
        SECONDARY,
    )
    add_icon_card(
        slide,
        8.0,
        4.6,
        4.8,
        1.45,
        "Better Efficiency",
        "Reduces manual effort while scaling to many users and items.",
        "EF",
        ACCENT,
    )


def add_slide_3_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "3. Problem Statement", 3)

    problem_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.4), Inches(11.9), Inches(4.9)
    )
    problem_box.fill.solid()
    problem_box.fill.fore_color.rgb = WHITE
    problem_box.line.color.rgb = RGBColor(0xD7, 0xDF, 0xE8)

    bullets = [
        "Traditional auctions require physical presence, limiting participation and reach.",
        "Managing bidder records and bid updates manually is difficult and error-prone.",
        "Manual operations often cause delays, inconsistency, and operational overhead.",
        "Lack of transparency can reduce trust among sellers, buyers, and organizers.",
    ]
    add_bullets(slide, 1.3, 2.0, 10.9, 3.4, bullets, font_size=19)

    for i, code in enumerate(["PH", "MG", "ER", "TR"]):
        circ = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95 + i * 3.0), Inches(5.55), Inches(0.62), Inches(0.62)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT if i % 2 == 0 else SECONDARY
        circ.line.fill.background()
        set_box_text(circ, code, font_name="Montserrat", size=11, bold=True, color=WHITE)


def add_slide_4_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "4. Objectives", 4)

    objectives = [
        "Develop a secure and reliable auction platform.",
        "Enable user registration and authenticated login.",
        "Allow sellers to list and manage auction items.",
        "Allow buyers to place bids with real-time updates.",
        "Automatically declare the highest bidder as the winner.",
    ]

    y = 1.45
    for idx, item in enumerate(objectives, start=1):
        row = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(11.3), Inches(0.86)
        )
        row.fill.solid()
        row.fill.fore_color.rgb = WHITE
        row.line.color.rgb = RGBColor(0xD7, 0xDF, 0xE8)

        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.2), Inches(y + 0.15), Inches(0.56), Inches(0.56))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PRIMARY if idx % 2 else SECONDARY
        badge.line.fill.background()
        set_box_text(badge, str(idx), font_name="Montserrat", size=12, bold=True, color=WHITE)

        text = slide.shapes.add_textbox(Inches(1.95), Inches(y + 0.18), Inches(9.8), Inches(0.45))
        set_box_text(text, item, font_name="Calibri", size=18, color=TEXT_DARK, align=PP_ALIGN.LEFT)
        y += 1.03


def add_slide_5_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "5. System Requirements", 5)

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(5.95), Inches(5.6)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = SECONDARY

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.58), Inches(1.3), Inches(5.95), Inches(5.6)
    )
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = ACCENT

    set_box_text(left, "Functional Requirements", font_name="Poppins", size=20, bold=True, color=PRIMARY)
    set_box_text(right, "Non-Functional Requirements", font_name="Poppins", size=20, bold=True, color=PRIMARY)

    add_bullets(
        slide,
        1.1,
        2.1,
        5.3,
        4.3,
        [
            "User registration and profile creation",
            "Secure user login and authentication",
            "Sellers can list and update auction items",
            "Buyers can place bids in active auctions",
            "Auction timer to start/end each listing",
        ],
        font_size=15,
    )

    add_bullets(
        slide,
        6.9,
        2.1,
        5.3,
        4.3,
        [
            "Data security and privacy protection",
            "Reliable operation with minimal downtime",
            "Fast response time during peak bidding",
            "User-friendly and intuitive interface",
            "Scalable design for future growth",
        ],
        font_size=15,
    )


def add_slide_6_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "6. System Architecture", 6)

    add_bullets(
        slide,
        0.85,
        1.0,
        12.1,
        0.8,
        ["Main components: Frontend UI, Backend Server, Auction Management, Bidding System, Database, and Payment Module."],
        font_size=14,
    )

    components = [
        ("User Interface\n(Frontend)", 0.8, 2.0, PRIMARY),
        ("Backend\nServer", 3.2, 2.0, SECONDARY),
        ("Auction\nManagement", 5.6, 2.0, PRIMARY),
        ("Bidding\nSystem", 8.0, 2.0, SECONDARY),
        ("Database", 4.2, 4.3, PRIMARY),
        ("Payment\nModule", 7.0, 4.3, ACCENT),
    ]

    for label, x, y, color in components:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.1), Inches(1.05)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        set_box_text(box, label, font_name="Calibri", size=13, bold=True, color=TEXT_DARK)

    add_arrow(slide, 2.95, 2.45)
    add_arrow(slide, 5.35, 2.45)
    add_arrow(slide, 7.75, 2.45)

    down1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(5.0), Inches(3.15), Inches(0.35), Inches(0.65))
    down1.fill.solid()
    down1.fill.fore_color.rgb = SECONDARY
    down1.line.fill.background()

    down2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(8.1), Inches(3.15), Inches(0.35), Inches(0.65))
    down2.fill.solid()
    down2.fill.fore_color.rgb = SECONDARY
    down2.line.fill.background()


def add_slide_7_use_case(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "7. Use Case Diagram", 7)

    boundary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.35), Inches(6.95), Inches(5.35)
    )
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(0xF9, 0xFB, 0xFF)
    boundary.line.color.rgb = SECONDARY

    boundary_title = slide.shapes.add_textbox(Inches(3.45), Inches(1.5), Inches(6.45), Inches(0.35))
    set_box_text(boundary_title, "Auction Management System", font_name="Poppins", size=16, bold=True, color=PRIMARY)

    actors = [("Admin", 0.65, 2.1), ("Seller", 0.65, 3.35), ("Buyer", 0.65, 4.6)]
    for name, x, y in actors:
        actor = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.75))
        actor.fill.solid()
        actor.fill.fore_color.rgb = WHITE
        actor.line.color.rgb = PRIMARY
        set_box_text(actor, name, font_name="Calibri", size=15, bold=True, color=PRIMARY)

    use_cases = [
        ("Register", 3.7, 2.0),
        ("Login", 6.95, 2.0),
        ("List Item", 3.7, 3.3),
        ("Place Bid", 6.95, 3.3),
        ("Manage Auctions", 3.7, 4.6),
        ("Declare Winner", 6.95, 4.6),
    ]

    for uc, x, y in use_cases:
        oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(2.45), Inches(0.78))
        oval.fill.solid()
        oval.fill.fore_color.rgb = WHITE
        oval.line.color.rgb = SECONDARY
        set_box_text(oval, uc, font_name="Calibri", size=13, bold=True, color=TEXT_DARK)

    # Association lines
    for y in [2.35, 3.6, 4.85]:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Inches(2.45), Inches(y), Inches(1.15), Inches(0.0))
        line.line.color.rgb = PRIMARY


def add_slide_8_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "8. System Workflow", 8)

    steps = [
        "1. User\nRegistration",
        "2. Login",
        "3. Seller\nLists Item",
        "4. Auction\nBegins",
        "5. Buyers\nPlace Bids",
        "6. Highest Bid\nUpdated",
        "7. Auction\nEnds",
        "8. Winner\nDeclared",
    ]

    x = 0.6
    y = 2.6
    box_w = 1.45
    box_h = 1.2

    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY if i % 2 == 0 else SECONDARY
        set_box_text(box, step, font_name="Calibri", size=11, bold=True, color=TEXT_DARK)

        if i < len(steps) - 1:
            add_arrow(slide, x + box_w + 0.03, y + 0.49, w=0.32, h=0.22)
        x += 1.62

    caption = slide.shapes.add_textbox(Inches(0.9), Inches(4.35), Inches(11.7), Inches(1.1))
    set_box_text(
        caption,
        "The workflow ensures a fair and automated process from user onboarding to winner declaration.",
        font_name="Calibri",
        size=16,
        color=PRIMARY,
        align=PP_ALIGN.LEFT,
    )


def add_slide_9_advantages(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "9. Advantages", 9)

    advantages = [
        ("Easy Access", "Users can join auctions anytime and from anywhere.", PRIMARY),
        ("Transparent Bidding", "Live bid visibility builds trust and fairness.", SECONDARY),
        ("Reduced Manual Work", "Automation lowers administrative effort.", ACCENT),
        ("Time Efficient", "Rapid updates speed up complete auction cycles.", PRIMARY),
        ("Better Management", "Centralized records simplify tracking and reporting.", SECONDARY),
    ]

    positions = [(0.85, 1.6), (4.45, 1.6), (8.05, 1.6), (2.65, 4.0), (6.25, 4.0)]

    for (title, desc, color), (x, y) in zip(advantages, positions):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.15), Inches(2.05))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color

        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.13), Inches(y + 0.13), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()

        tbox = slide.shapes.add_textbox(Inches(x + 0.68), Inches(y + 0.12), Inches(2.35), Inches(0.35))
        set_box_text(tbox, title, font_name="Poppins", size=14, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)

        dbox = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.62), Inches(2.82), Inches(1.25))
        set_box_text(dbox, desc, font_name="Calibri", size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT)


def add_slide_10_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_header(slide, "10. Conclusion and Future Scope", 10)

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.4), Inches(5.9), Inches(4.95)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = SECONDARY

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(1.4), Inches(5.9), Inches(4.95)
    )
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = ACCENT

    lt = slide.shapes.add_textbox(Inches(1.15), Inches(1.65), Inches(5.3), Inches(0.4))
    set_box_text(lt, "Conclusion", font_name="Poppins", size=20, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)

    add_bullets(
        slide,
        1.1,
        2.15,
        5.2,
        3.8,
        [
            "The system digitalizes auction operations through a structured SDLC approach.",
            "It increases efficiency, improves transparency, and reduces manual errors.",
            "Automated bid handling and winner declaration ensure fairness and reliability.",
        ],
        font_size=15,
    )

    rt = slide.shapes.add_textbox(Inches(6.85), Inches(1.65), Inches(5.2), Inches(0.4))
    set_box_text(rt, "Future Scope", font_name="Poppins", size=20, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)

    add_bullets(
        slide,
        6.8,
        2.15,
        5.1,
        3.8,
        [
            "Develop a dedicated mobile application.",
            "Integrate AI-based fraud detection mechanisms.",
            "Provide real-time notifications and alerts.",
            "Explore blockchain-based secure auction models.",
        ],
        font_size=15,
    )

    thanks = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(6.55), Inches(3.6), Inches(0.62)
    )
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = PRIMARY
    thanks.line.fill.background()
    set_box_text(thanks, "Thank You", font_name="Poppins", size=20, bold=True, color=WHITE)


def build_presentation(output_path="Auction_Management_System_SDLC_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)
    add_slide_2_intro(prs)
    add_slide_3_problem(prs)
    add_slide_4_objectives(prs)
    add_slide_5_requirements(prs)
    add_slide_6_architecture(prs)
    add_slide_7_use_case(prs)
    add_slide_8_workflow(prs)
    add_slide_9_advantages(prs)
    add_slide_10_conclusion(prs)

    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    build_presentation()
